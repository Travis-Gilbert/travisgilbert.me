"""
Multi-format file extraction for CommonPlace Object ingestion.

Supported formats:
  PDF   - spacy-layout / pypdf (existing, delegated to pdf_ingestion.py)
  DOCX  - python-docx (structure-aware: headings, paragraphs, tables)
  XLSX  - openpyxl (sheet names, cell data as structured text)
  PPTX  - python-pptx (slide titles + body text)
  Images - Pillow for metadata, pytesseract for OCR, optional scikit-image
  TXT/MD/CSV/HTML - direct text read
  Other - empty result with filename as title

Never raises. Returns:
  {
    'title': str,
    'body': str,
    'author': str,
    'sections': list[dict],   # [{heading, body}]
    'metadata': dict,          # format-specific metadata
    'char_count': int,
    'method': str,             # 'docx' | 'xlsx' | 'pptx' | 'image_ocr' | ...
    'thumbnails': list[str],   # base64-encoded PNG thumbnails
  }
"""

import logging
import mimetypes
from pathlib import Path
from typing import Any

from .pdf_ingestion import extract_pdf_text

logger = logging.getLogger(__name__)

CODE_EXTENSIONS = {'.py', '.js', '.ts', '.tsx', '.jsx'}


def extract_file_content(file_bytes: bytes, filename: str, mime_type: str = '') -> dict[str, Any]:
    """
    Route file to the appropriate extractor based on MIME type or extension.
    """
    if not mime_type:
        mime_type, _ = mimetypes.guess_type(filename)
        mime_type = mime_type or ''

    ext = Path(filename).suffix.lower()

    result: dict[str, Any]

    # PDF
    if mime_type == 'application/pdf' or ext == '.pdf':
        result = extract_pdf_text(file_bytes)

    # DOCX
    elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' or ext in ('.docx', '.doc'):
        result = _extract_docx(file_bytes)

    # XLSX
    elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' or ext in ('.xlsx', '.xls'):
        result = _extract_xlsx(file_bytes)

    # PPTX
    elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' or ext in ('.pptx', '.ppt'):
        result = _extract_pptx(file_bytes)

    # Code files
    elif ext in CODE_EXTENSIONS:
        result = _extract_code(file_bytes, filename, ext)

    # Images
    elif mime_type.startswith('image/') or ext in ('.jpg', '.jpeg', '.png', '.gif', '.webp', '.tiff', '.bmp'):
        result = _extract_image(file_bytes, filename, mime_type)

    # Plain text family
    elif mime_type.startswith('text/') or ext in ('.txt', '.md', '.csv', '.json', '.xml', '.yaml', '.yml', '.toml', '.html', '.htm'):
        result = _extract_text(file_bytes, ext)

    else:
        logger.warning('Unsupported file type: %s (%s)', filename, mime_type)
        result = _empty_result('unsupported')
        result['title'] = Path(filename).stem

    return _with_common_metadata(
        result=result,
        filename=filename,
        mime_type=mime_type,
        extension=ext,
        file_size=len(file_bytes),
    )


def _empty_result(method: str) -> dict[str, Any]:
    return {
        'title': '',
        'body': '',
        'author': '',
        'sections': [],
        'metadata': {},
        'char_count': 0,
        'method': method,
        'thumbnails': [],
    }


def _with_common_metadata(
    result: dict[str, Any],
    filename: str,
    mime_type: str,
    extension: str,
    file_size: int,
) -> dict[str, Any]:
    metadata = result.get('metadata')
    if not isinstance(metadata, dict):
        metadata = {}

    metadata.setdefault('filename', filename)
    metadata.setdefault('mime_type', mime_type or 'application/octet-stream')
    metadata.setdefault('extension', extension)
    metadata.setdefault('file_size', file_size)
    metadata.setdefault('parser', result.get('method', 'unknown'))

    result['metadata'] = metadata
    result['char_count'] = int(result.get('char_count') or len(result.get('body') or ''))
    return result


# ─── DOCX ────────────────────────────────────────────────────────────────────


def _extract_docx(file_bytes: bytes) -> dict[str, Any]:
    """
    Extract structured content from a DOCX file.
    Headings become section delimiters. Tables render as pipe-separated rows.
    Core properties supply title and author metadata.
    """
    try:
        import io
        from docx import Document

        doc = Document(io.BytesIO(file_bytes))
        props = doc.core_properties
        title = (props.title or '').strip()
        author = (props.author or '').strip()

        sections = []
        current_heading = ''
        current_body: list[str] = []
        all_text: list[str] = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            style_name = (para.style.name or '').lower()

            if 'heading' in style_name:
                if current_heading or current_body:
                    sections.append({
                        'heading': current_heading,
                        'body': '\n'.join(current_body),
                    })
                    current_body = []
                current_heading = text
                if not title:
                    title = text
            else:
                current_body.append(text)

            all_text.append(text)

        if current_heading or current_body:
            sections.append({
                'heading': current_heading,
                'body': '\n'.join(current_body),
            })

        # Tables as pipe-delimited text
        table_text: list[str] = []
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(' | '.join(cells))
            if rows:
                table_text.append('\n'.join(rows))

        body = '\n\n'.join(all_text)
        if table_text:
            body += '\n\n---\n\n' + '\n\n'.join(table_text)

        if not title and all_text:
            title = all_text[0][:80]

        return {
            'title': title,
            'body': body,
            'author': author,
            'sections': sections,
            'metadata': {
                'section_count': len(doc.sections),
                'table_count': len(doc.tables),
                'created': str(props.created) if props.created else '',
                'modified': str(props.modified) if props.modified else '',
            },
            'char_count': len(body),
            'method': 'docx',
            'thumbnails': [],
        }
    except Exception as exc:
        logger.warning('DOCX extraction failed: %s', exc)
        return _empty_result('docx_failed')


# ─── XLSX ────────────────────────────────────────────────────────────────────


def _extract_xlsx(file_bytes: bytes) -> dict[str, Any]:
    """
    Extract structured content from an XLSX spreadsheet.
    Each sheet becomes a section. Cell values are joined as tab-separated rows.
    """
    try:
        import io
        import openpyxl

        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        sections = []
        all_text: list[str] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(cell) if cell is not None else '' for cell in row]
                if any(cells):
                    rows.append('\t'.join(cells))

            body = '\n'.join(rows)
            sections.append({'heading': sheet_name, 'body': body})
            all_text.append(f'# {sheet_name}\n{body}')

        full_body = '\n\n'.join(all_text)
        title = wb.sheetnames[0] if wb.sheetnames else ''

        return {
            'title': title,
            'body': full_body,
            'author': '',
            'sections': sections,
            'metadata': {'sheet_count': len(wb.sheetnames)},
            'char_count': len(full_body),
            'method': 'xlsx',
            'thumbnails': [],
        }
    except Exception as exc:
        logger.warning('XLSX extraction failed: %s', exc)
        return _empty_result('xlsx_failed')


# ─── PPTX ────────────────────────────────────────────────────────────────────


def _extract_pptx(file_bytes: bytes) -> dict[str, Any]:
    """
    Extract text from a PowerPoint presentation.
    Each slide becomes a section with its title and body text.
    """
    try:
        import io
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_bytes))
        sections = []
        all_text: list[str] = []
        presentation_title = ''

        for i, slide in enumerate(prs.slides, 1):
            slide_title = ''
            slide_body_parts: list[str] = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue

                # Title placeholder (placeholder type 15 or name contains 'title')
                is_title = (
                    hasattr(shape, 'placeholder_format')
                    and shape.placeholder_format is not None
                    and shape.placeholder_format.idx in (0, 1)
                )
                if is_title and not slide_title:
                    slide_title = text
                    if i == 1 and not presentation_title:
                        presentation_title = text
                else:
                    slide_body_parts.append(text)

            body = '\n'.join(slide_body_parts)
            heading = slide_title or f'Slide {i}'
            sections.append({'heading': heading, 'body': body})
            chunk = f'# {heading}\n{body}' if body else f'# {heading}'
            all_text.append(chunk)

        full_body = '\n\n'.join(all_text)

        return {
            'title': presentation_title,
            'body': full_body,
            'author': '',
            'sections': sections,
            'metadata': {'slide_count': len(prs.slides)},
            'char_count': len(full_body),
            'method': 'pptx',
            'thumbnails': [],
        }
    except Exception as exc:
        logger.warning('PPTX extraction failed: %s', exc)
        return _empty_result('pptx_failed')


# ─── Images ──────────────────────────────────────────────────────────────────


def _extract_image(file_bytes: bytes, filename: str, mime_type: str) -> dict[str, Any]:
    """
    Image ingestion with two tiers:
    Tier 1: Pillow for metadata and thumbnail + pytesseract OCR (always attempted)
    Tier 2: scikit-image dominant color extraction (when installed)

    SAM-2 segmentation is intentionally async-only and not called here.
    """
    result = _empty_result('image')
    result['title'] = Path(filename).stem

    try:
        import base64
        import io
        from PIL import Image
        from PIL.ExifTags import TAGS

        img = Image.open(io.BytesIO(file_bytes))
        width, height = img.size

        # EXIF metadata
        exif_data: dict[str, str] = {}
        try:
            raw_exif = img.getexif()
            for tag_id, value in raw_exif.items():
                tag_name = TAGS.get(tag_id, str(tag_id))
                exif_data[tag_name] = str(value)
        except Exception:
            pass

        result['metadata'] = {
            'width': width,
            'height': height,
            'format': img.format or mime_type.split('/')[-1],
            'mode': img.mode,
            'exif': exif_data,
            'sam_available': False,
        }

        # Thumbnail (max 200px, PNG)
        thumb = img.copy()
        thumb.thumbnail((200, 200))
        thumb_buf = io.BytesIO()
        thumb_rgb = thumb.convert('RGB')
        thumb_rgb.save(thumb_buf, format='PNG')
        result['thumbnails'] = [base64.b64encode(thumb_buf.getvalue()).decode()]

        # OCR with pytesseract
        try:
            import pytesseract
            ocr_text = pytesseract.image_to_string(img).strip()
            if ocr_text and len(ocr_text) > 10:
                result['body'] = ocr_text
                result['char_count'] = len(ocr_text)
                result['method'] = 'image_ocr'
        except ImportError:
            logger.debug('pytesseract not installed; skipping OCR')
        except Exception as ocr_exc:
            logger.warning('OCR failed: %s', ocr_exc)

        # Tier 2: dominant colors via scikit-image + sklearn
        result = _try_dominant_colors(result, img)

    except Exception as exc:
        logger.warning('Image extraction failed: %s', exc)

    return result


def _try_dominant_colors(result: dict[str, Any], img) -> dict[str, Any]:
    """
    Optional Tier 2: extract dominant colors with KMeans clustering.
    Runs only when scikit-learn is installed. Lightweight, synchronous.
    """
    try:
        import numpy as np
        from sklearn.cluster import KMeans

        img_rgb = img.convert('RGB')
        pixels = np.array(img_rgb).reshape(-1, 3).astype(float)

        # Subsample large images for speed
        if len(pixels) > 2000:
            indices = np.random.default_rng(42).choice(len(pixels), 2000, replace=False)
            pixels = pixels[indices]

        kmeans = KMeans(n_clusters=4, n_init=3, random_state=42)
        kmeans.fit(pixels)

        dominant = [
            '#{:02x}{:02x}{:02x}'.format(int(c[0]), int(c[1]), int(c[2]))
            for c in kmeans.cluster_centers_
        ]
        result['metadata']['dominant_colors'] = dominant
        if result['method'] == 'image':
            result['method'] = 'image_color'

    except ImportError:
        pass
    except Exception as exc:
        logger.debug('Dominant color extraction skipped: %s', exc)

    return result


# ─── Code files (tree-sitter + fallback) ────────────────────────────────────


def _extract_code(file_bytes: bytes, filename: str, ext: str) -> dict[str, Any]:
    """
    Structured code extraction for .py/.js/.ts/.tsx/.jsx.

    Priority:
    1) tree-sitter AST walk (imports/functions/classes/comments)
    2) safe fallback parser (ast for Python, regex for JS/TS)
    """
    source = file_bytes.decode('utf-8', errors='replace')
    title = Path(filename).stem

    structured = _extract_code_with_tree_sitter(source, ext)
    method = 'code_tree_sitter'
    if structured is None:
        structured = _extract_code_with_fallback(source, ext)
        method = 'code_fallback'

    imports = structured.get('imports', [])
    functions = structured.get('functions', [])
    classes = structured.get('classes', [])
    comment_summary = structured.get('comment_summary', '')
    parser_name = structured.get('parser', method)
    docstrings = structured.get('docstrings', [])

    sections: list[dict[str, str]] = []
    if imports:
        sections.append({'heading': 'Imports', 'body': '\n'.join(imports[:80])})
    if classes:
        sections.append({'heading': 'Classes', 'body': '\n'.join(classes[:80])})
    if functions:
        sections.append({'heading': 'Functions', 'body': '\n'.join(functions[:120])})
    if docstrings:
        sections.append({'heading': 'Docstrings', 'body': '\n\n'.join(docstrings[:12])})
    if comment_summary:
        sections.append({'heading': 'Comment Summary', 'body': comment_summary})

    body_chunks = [
        '# Code Overview',
        f'Language: {_language_label_from_extension(ext)}',
        f'Imports: {len(imports)}',
        f'Classes: {len(classes)}',
        f'Functions: {len(functions)}',
    ]
    if comment_summary:
        body_chunks.append('\n# Notes\n' + comment_summary)
    body_chunks.append('\n# Source\n' + source[:8000])
    body = '\n'.join(body_chunks)

    return {
        'title': title,
        'body': body,
        'author': '',
        'sections': sections,
        'metadata': {
            'language': _language_label_from_extension(ext),
            'import_count': len(imports),
            'class_count': len(classes),
            'function_count': len(functions),
            'docstring_count': len(docstrings),
            'comment_summary': comment_summary,
            'parser': parser_name,
        },
        'char_count': len(body),
        'method': method,
        'thumbnails': [],
    }


def _extract_code_with_tree_sitter(source: str, ext: str) -> dict[str, Any] | None:
    try:
        from tree_sitter import Language, Parser
    except Exception:
        return None

    language_obj = _load_tree_sitter_language(ext)
    if language_obj is None:
        return None

    try:
        parser = Parser()
        if hasattr(parser, 'set_language'):
            parser.set_language(language_obj)
        else:
            parser = Parser(language_obj)
        tree = parser.parse(source.encode('utf-8', errors='ignore'))
    except Exception as exc:
        logger.debug('tree-sitter parse failed for %s: %s', ext, exc)
        return None

    imports: list[str] = []
    functions: list[str] = []
    classes: list[str] = []
    comments: list[str] = []

    import_types = {
        'import_statement',
        'import_from_statement',
        'import_declaration',
    }
    function_types = {
        'function_definition',
        'function_declaration',
        'method_definition',
        'generator_function_declaration',
        'arrow_function',
    }
    class_types = {'class_definition', 'class_declaration'}

    root = tree.root_node
    stack = [root]
    seen: set[tuple[str, int, int]] = set()

    while stack:
        node = stack.pop()
        for child in node.children:
            stack.append(child)

        snippet = source[node.start_byte:node.end_byte].strip()
        if not snippet:
            continue

        key = (node.type, node.start_byte, node.end_byte)
        if key in seen:
            continue
        seen.add(key)

        if node.type in import_types:
            imports.append(snippet.splitlines()[0][:220])
            continue

        if node.type in class_types:
            classes.append(_ts_named_signature(node, source))
            continue

        if node.type in function_types:
            sig = _ts_named_signature(node, source)
            if sig not in functions:
                functions.append(sig)
            continue

        if node.type == 'comment':
            comments.append(snippet)

    if not any([imports, functions, classes, comments]):
        return None

    return {
        'imports': imports,
        'functions': functions,
        'classes': classes,
        'docstrings': [],
        'comment_summary': _summarize_comments(comments),
        'parser': 'tree_sitter',
    }


def _load_tree_sitter_language(ext: str):
    try:
        from tree_sitter import Language
    except Exception:
        return None

    module = None
    if ext == '.py':
        try:
            import tree_sitter_python as module
        except Exception:
            return None
    elif ext in {'.js', '.jsx'}:
        try:
            import tree_sitter_javascript as module
        except Exception:
            return None
    elif ext in {'.ts', '.tsx'}:
        try:
            import tree_sitter_typescript as module
        except Exception:
            return None
    else:
        return None

    try:
        lang_func = getattr(module, 'language')
        raw_language = lang_func()
        try:
            return Language(raw_language)
        except Exception:
            return raw_language
    except Exception:
        return None


def _ts_named_signature(node, source: str) -> str:
    name_node = None
    try:
        name_node = node.child_by_field_name('name')
    except Exception:
        name_node = None

    if name_node is not None:
        name = source[name_node.start_byte:name_node.end_byte].strip()
        return f'{node.type}: {name}'

    snippet = source[node.start_byte:node.end_byte].strip().splitlines()[0]
    return f'{node.type}: {snippet[:180]}'


def _extract_code_with_fallback(source: str, ext: str) -> dict[str, Any]:
    if ext == '.py':
        return _extract_python_fallback(source)
    return _extract_js_like_fallback(source)


def _extract_python_fallback(source: str) -> dict[str, Any]:
    import ast

    imports: list[str] = []
    functions: list[str] = []
    classes: list[str] = []
    docstrings: list[str] = []

    try:
        tree = ast.parse(source)
    except Exception:
        tree = None

    if tree is not None:
        for node in ast.walk(tree):
            if isinstance(node, ast.Import):
                imports.extend(alias.name for alias in node.names)
            elif isinstance(node, ast.ImportFrom):
                mod = node.module or ''
                imported = ', '.join(alias.name for alias in node.names)
                imports.append(f'from {mod} import {imported}')
            elif isinstance(node, ast.FunctionDef):
                functions.append(node.name)
                ds = ast.get_docstring(node)
                if ds:
                    docstrings.append(f'{node.name}: {ds[:280]}')
            elif isinstance(node, ast.AsyncFunctionDef):
                functions.append(f'async {node.name}')
                ds = ast.get_docstring(node)
                if ds:
                    docstrings.append(f'{node.name}: {ds[:280]}')
            elif isinstance(node, ast.ClassDef):
                classes.append(node.name)
                ds = ast.get_docstring(node)
                if ds:
                    docstrings.append(f'{node.name}: {ds[:280]}')

    comments = []
    for line in source.splitlines():
        stripped = line.strip()
        if stripped.startswith('#'):
            comments.append(stripped)

    return {
        'imports': imports,
        'functions': functions,
        'classes': classes,
        'docstrings': docstrings,
        'comment_summary': _summarize_comments(comments),
        'parser': 'python_ast',
    }


def _extract_js_like_fallback(source: str) -> dict[str, Any]:
    import re

    imports = re.findall(r'^(?:import\\s+.+|const\\s+.+?=\\s*require\\(.+\\))', source, flags=re.MULTILINE)
    function_defs = re.findall(
        r'^(?:export\\s+)?(?:async\\s+)?function\\s+([A-Za-z_$][\\w$]*)',
        source,
        flags=re.MULTILINE,
    )
    class_defs = re.findall(
        r'^(?:export\\s+)?class\\s+([A-Za-z_$][\\w$]*)',
        source,
        flags=re.MULTILINE,
    )
    comments = re.findall(r'//.*?$|/\\*[\\s\\S]*?\\*/', source, flags=re.MULTILINE)

    return {
        'imports': imports,
        'functions': function_defs,
        'classes': class_defs,
        'docstrings': [],
        'comment_summary': _summarize_comments(comments),
        'parser': 'regex',
    }


def _language_label_from_extension(ext: str) -> str:
    mapping = {
        '.py': 'python',
        '.js': 'javascript',
        '.jsx': 'javascript',
        '.ts': 'typescript',
        '.tsx': 'typescript',
    }
    return mapping.get(ext, ext.lstrip('.') or 'unknown')


def _summarize_comments(comments: list[str]) -> str:
    cleaned: list[str] = []
    for comment in comments[:12]:
        line = (
            comment.replace('/*', '')
            .replace('*/', '')
            .replace('//', '')
            .replace('#', '')
            .strip()
        )
        if line:
            cleaned.append(line)
    return '\n'.join(cleaned[:6])


# ─── Plain text family ───────────────────────────────────────────────────────


def _extract_text(file_bytes: bytes, ext: str) -> dict[str, Any]:
    """
    Read plain-text files (TXT, MD, CSV, JSON, YAML, HTML, etc.).
    For HTML, strips tags with BeautifulSoup if available; falls back to raw decode.
    """
    try:
        raw = file_bytes.decode('utf-8', errors='replace')

        if ext in ('.html', '.htm'):
            raw = _strip_html(raw)

        lines = [line.strip() for line in raw.splitlines() if line.strip()]
        title = lines[0][:80] if lines else ''
        body = '\n'.join(lines)

        return {
            'title': title,
            'body': body,
            'author': '',
            'sections': [],
            'metadata': {'extension': ext},
            'char_count': len(body),
            'method': 'text',
            'thumbnails': [],
        }
    except Exception as exc:
        logger.warning('Text extraction failed: %s', exc)
        return _empty_result('text_failed')


def _strip_html(raw: str) -> str:
    """Strip HTML tags. Uses BeautifulSoup when available, regex fallback."""
    try:
        from bs4 import BeautifulSoup
        return BeautifulSoup(raw, 'html.parser').get_text(separator='\n')
    except ImportError:
        import re
        return re.sub(r'<[^>]+>', '', raw)
